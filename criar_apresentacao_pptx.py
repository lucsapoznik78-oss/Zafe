#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Gerador de apresentação Zafe em PowerPoint (.pptx)
Requer: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Cores
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(30, 30, 30)
GREEN = RGBColor(134, 239, 172)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 200)
RED = RGBColor(239, 68, 68)
YELLOW = RGBColor(234, 179, 8)

def add_slide(prs, layout_index=6):  # 6 = Blank layout
    return prs.slides.add_slide(prs.slide_layouts[layout_index])

def add_textbox(slide, text, left, top, width, height, 
                font_size=18, bold=False, color=WHITE, 
                bg_color=None, alignment=PP_ALIGN.LEFT):
    """Adiciona uma caixa de texto ao slide"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    
    if bg_color:
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = bg_color
    
    return textbox

def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=16, color=WHITE, bullet=False, alignment=PP_ALIGN.LEFT):
    """Adiciona texto com múltiplas linhas"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.alignment = alignment
        if bullet:
            p.level = 0
    
    return textbox

def set_slide_background(slide, color=DARK_GRAY):
    """Define a cor de fundo do slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 1280px equivalent
    prs.slide_height = Inches(7.5)     # 720px equivalent
    
    # ========== SLIDE 1: CAPA ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    # Logo ZAFE
    add_textbox(slide, "ZAFE", 1, 1.5, 11, 2, 
                font_size=96, bold=True, color=GREEN, 
                alignment=PP_ALIGN.CENTER)
    
    # Subtítulo
    add_textbox(slide, "A Liga de Previsões do Brasil", 1, 3.5, 11, 1,
                font_size=32, bold=False, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    
    # Créditos
    add_textbox(slide, "FEITO POR LUC SAPOZNIK E LORENZO FRAGALI", 
                1, 6.5, 11, 0.5,
                font_size=18, bold=True, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    
    # NOTAS DO PALESTRANTE
    slide.notes_slide.notes_text_frame.text = \
        "Zafe: liga de previsões (não é aposta, não é cassino). " \
        "Três pilares: Econômico, Liga, Privadas. " \
        "Moeda virtual: Zafes. Prêmios reais via PIX no concurso. " \
        "Apresentação: Luc Sapoznik e Lorenzo Fragali."
    
    # ========== SLIDE 2: QUEM SOMOS ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Quem Somos", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "A Zafe é a plataforma onde você compete prevendo o que vai acontecer no Brasil e no mundo.", 
                1, 1.8, 11, 0.8,
                font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # O que somos
    add_textbox(slide, "O que somos", 1, 2.8, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_esquerda = [
        "• 100% Legal (Lei 5.768/71 + CMN 5.298)",
        "• Sem apostas, sem depósitos de dinheiro real",
        "• Competição de habilidade analítica",
        "• Moeda virtual: Zafes (não convertível em dinheiro)",
        "• Prêmios reais via PIX no concurso"
    ]
    add_multiline_text(slide, items_esquerda, 1, 3.4, 5.5, 2.5,
                       font_size=16, color=WHITE, bullet=False)
    
    # O que não somos
    add_textbox(slide, "O que NÃO somos", 7, 2.8, 5, 0.5,
                font_size=24, bold=True, color=RED)
    
    items_direita = [
        "• Não somos casa de apostas",
        "• Não somos mercado preditivo (Polymarket)",
        "• Não operamos com dinheiro real (exceto prêmios)",
        "• Não temos odds ou 'casa' contra o usuário",
        "• Não existe depósito ou saque de Reais"
    ]
    add_multiline_text(slide, items_direita, 7, 3.4, 5.5, 2.5,
                       font_size=16, color=WHITE, bullet=False)
    
    slide.notes_slide.notes_text_frame.text = \
        "Quem somos: plataforma de competição de previsões (100% legal). " \
        "Não somos casa de apostas, não operamos com dinheiro real. " \
        "Moeda virtual: Zafes. Prêmios reais via PIX no concurso (Lei 5.768/71). " \
        "Diferencial: habilidade, não sorte."
    
    # ========== SLIDE 3: POR QUE LIGA ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Por que a Zafe existe?", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "O Problema", 1, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    problemas = [
        "• Brasileiro é proibido de prever eventos esportivos (CMN 5.298/2026)",
        "• Casas de aposta exploram dependência química (dopamina)",
        "• Usuário perde dinheiro real sem hiberitação analítica",
        "• Falta de plataforma nacional de competição de previsões"
    ]
    add_multiline_text(slide, problemas, 1, 2.6, 5.5, 2.5,
                       font_size=16, color=WHITE)
    
    add_textbox(slide, "A Solução Zafe", 7, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    solucoes = [
        "• Posicionamento legal: concurso de previsões (Lei 5.768/71)",
        "• Sem dinheiro real em jogo (Zafes virtuais)",
        "• Foco em habilidade analítica e calibração",
        "• Curadoria de informação (Premium) via IA"
    ]
    add_multiline_text(slide, solucoes, 7, 2.6, 5.5, 2.5,
                       font_size=16, color=WHITE)
    
    slide.notes_slide.notes_text_frame.text = \
        "Problema: brasileiro proibido de prever eventos (CMN 5.298/2026). " \
        "Casas de aposta exploram dependência. Solução: Zafe como concurso legal. " \
        "Sem dinheiro real, foco em habilidade e curadoria via IA (Premium)."
    
    # ========== SLIDE 4: ECONOMICO ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Zafe Econômico", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "Palpite sobre indicadores econômicos", 1, 1.8, 11, 0.6,
                font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Esquerda: Como funciona
    add_textbox(slide, "Como Funciona", 1, 2.6, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_eco_1 = [
        "• Criação: Apenas admin/sistema (não usuários)",
        "• Moeda: Z$ Virtual (não convertível em dinheiro real)",
        "• Eventos: Selic, IPCA, Dólar, Bitcoin, Ibovespa",
        "• Mercado Secundário: Order book com matching FIFO"
    ]
    add_multiline_text(slide, items_eco_1, 1, 3.2, 5.5, 2,
                       font_size=16, color=WHITE)
    
    # Direita: Categorias
    add_textbox(slide, "Categorias Permitidas", 7, 2.6, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_eco_2 = [
        "• Macro Brasil: Selic, IPCA, PIB, desemprego",
        "• Mercados: Ibovespa, Dólar, Euro, Ouro",
        "• Cripto: Bitcoin, Ethereum (janelas curtas)",
        "• Indicadores Globais: Fed, BCE, CPI americano"
    ]
    add_multiline_text(slide, items_eco_2, 7, 3.2, 5.5, 2,
                       font_size=16, color=WHITE)
    
    slide.notes_slide.notes_text_frame.text = \
        "Econômico: palpites sobre indicadores econômicos (Z$ virtual). " \
        "Apenas admin cria eventos (Selic, IPCA, Dólar, Bitcoin, Ibovespa). " \
        "Mercado secundário com order book FIFO. Lei 5.768/71."
    
    # ========== SLIDE 5: LIGA ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Zafe Liga (O Carro-Chefe)", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "A liga onde você compete prevendo o que vai acontecer", 
                1, 1.8, 11, 0.6,
                font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # Esquerda
    add_textbox(slide, "Como Funciona", 1, 2.6, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_liga_1 = [
        "• Criação: Usuários e Admin (aprovação necessária)",
        "• Moeda: Z$ Virtual (mesma da Liga e Privadas)",
        "• Universo: Esporte, Política, BBB, Oscar, Tecnologia",
        "• Banner do Concurso ativo com inscrição grátis"
    ]
    add_multiline_text(slide, items_liga_1, 1, 3.2, 5.5, 2,
                       font_size=16, color=WHITE)
    
    # Direita
    add_textbox(slide, "Diferenciais", 7, 2.6, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_liga_2 = [
        "• Ranking por Brier Score (acurácia agregada)",
        "• Mínimo 30 previsões/mês para qualificação",
        "• Diversidade: pelo menos 3 categorias diferentes",
        "• Trending Feed: eventos em alta (+Z$ / 2h)"
    ]
    add_multiline_text(slide, items_liga_2, 7, 3.2, 5.5, 2,
                       font_size=16, color=WHITE)
    
    slide.notes_slide.notes_text_frame.text = \
        "Liga: carro-chefe da Zafe (esporte, política, BBB, Oscar, tecnologia). " \
        "Usuários criam eventos (aprovação admin). Moeda: Z$ virtual. " \
        "Ranking por Brier Score, mínimo 30 previsões/mês, 3 categorias. " \
        "Banner do concurso com inscrição grátis."
    
    # ========== SLIDE 6: CONCURSO ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Concurso (Sobreposto à Liga)", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    # Esquerda
    add_textbox(slide, "Como Funciona", 1, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_conc_1 = [
        "• Inscrição: Gratuita na fase beta (futuramente R$ 9,90)",
        "• Critério: Maior acurácia agregada (Brier Score)",
        "• Diversidade: Previsões em pelo menos 3 categorias",
        "• Prêmio: R$ 2.500 totais distribuídos via PIX",
        "• Recurso: Patrocínio de marcas (Coca-Cola, bancos)"
    ]
    add_multiline_text(slide, items_conc_1, 1, 2.6, 5.5, 2.5,
                       font_size=16, color=WHITE)
    
    add_textbox(slide, "Distribuição de Prêmios:", 1, 5.2, 5, 0.4,
                font_size=18, bold=True, color=GREEN)
    
    premios = [
        "1º lugar: R$ 1.000 | 2º lugar: R$ 500 | 3º lugar: R$ 250",
        "4º a 10º lugar: R$ 100 cada | 11º a 25º lugar: R$ 35 cada"
    ]
    add_multiline_text(slide, premios, 1, 5.6, 5.5, 1,
                       font_size=14, color=LIGHT_GRAY)
    
    # Direita
    add_textbox(slide, "Base Legal (Lei 5.768/71)", 7, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    add_textbox(slide, 
                "Art. 1º: Concurso de previsões exige autorização SECAP apenas se prêmios > 10% do valor arrecadado.",
                7, 2.6, 5.5, 0.8,
                font_size=14, color=WHITE)
    
    add_textbox(slide, "Dispensa SECAP", 7, 3.5, 5, 0.4,
                font_size=18, bold=True, color=GREEN)
    
    add_textbox(slide,
                "Concurso regular (R$ 2.500) tem patrocínio que cobre prêmios. Não há arrecadação líquida da Zafe, dispensando autorização.",
                7, 4, 5.5, 1,
                font_size=14, color=WHITE)
    
    add_textbox(slide, "Concurso Major: Trimestral R$ 10k-25k. Exige processo formal na SECAP, custo de 10% pago pelo patrocinador.",
                7, 5.2, 5.5, 1,
                font_size=14, color=YELLOW)
    
    slide.notes_slide.notes_text_frame.text = \
        "Concurso: R$ 2.500 via PIX (Lei 5.768/71). " \
        "Brier Score: maior acurácia agregada, mínimo 3 categorias. " \
        "Base legal: Dispensa SECAP (patrocínio cobre prêmios). " \
        "Concurso Major: R$ 10k-25k (exige SECAP, 10% patrocinador)."
    
    # ========== SLIDE 7: PRIVADAS ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Zafe Privadas (Bolão entre Amigos)", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    # Esquerda
    add_textbox(slide, "Como Funciona", 1, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_priv_1 = [
        "• Criação: Um usuário desafia outro(s) para um palpite",
        "• Mercado: Fechado (participantes definidos na criação)",
        "• Moeda: Z$ Virtual (não convertível em dinheiro real)",
        "• Resolução: Supermaioria 67% + Sistema de Juízes",
        "• Futuro: Converter Z$ em descontos reais via parcerias"
    ]
    add_multiline_text(slide, items_priv_1, 1, 2.6, 5.5, 2.5,
                       font_size=16, color=WHITE)
    
    # Direita
    add_textbox(slide, "4 Travas Legais (Obrigatórias)", 7, 2, 5, 0.5,
                font_size=24, bold=True, color=RED)
    
    travas = [
        "1. Apenas amigos confirmados (24h após aceite mútuo)",
        "2. Mercado fechado (sem novos participantes após criação)",
        "3. Sem revenda de posição (sem mercado secundário)",
        "4. Limite: 5.000 Z$/ano por par de usuários"
    ]
    add_multiline_text(slide, travas, 7, 2.6, 5.5, 2,
                       font_size=16, color=WHITE)
    
    add_textbox(slide, 
                "Sistema de Juízes: Se divergência na supermaioria (67%), ativa-se juízes imparciais para decidir.",
                7, 4.8, 5.5, 1,
                font_size=14, color=YELLOW)
    
    slide.notes_slide.notes_text_frame.text = \
        "Privadas: bolão entre amigos (4 travas legais obrigatórias). " \
        "Mercado fechado, sem revenda, limite 5.000 Z$/ano por par. " \
        "Futuro: converter Z$ em descontos reais via parcerias. " \
        "Sistema de Juízes: divergência na supermaioria (67%)."
    
    # ========== SLIDE 8: RESOLUCAO ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Como o Código Sabe quem Ganhou", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    # 4 camadas
    camadas = [
        ("Camada 1: API Fixa", "Banco Central\nYahoo Finance\nCoinGecko\nTSE, APIs Esportivas", 1, 2, GREEN),
        ("Camada 2: IA com Dupla Verificação", "2 checks independentes\nWeb Search\nConfiança ≥ 85%\nAmbos concordam", 4, 2, GREEN),
        ("Camada 3: Retry Automático", "Falhou?\nTenta novamente\na cada 2h\n(máximo 3x)", 7, 2, GREEN),
        ("Camada 4: Reembolso", "Ninguém perde\nse não resolver\nTodos recebem Z$\nde volta", 10, 2, RED)
    ]
    
    for titulo, conteudo, left, top, cor in camadas:
        add_textbox(slide, titulo, left, top, 2.5, 0.5,
                    font_size=18, bold=True, color=cor)
        add_textbox(slide, conteudo, left, top+0.6, 2.5, 1.5,
                    font_size=14, color=WHITE)
    
    add_textbox(slide, "Taxa de resolução automática: 87% | Reembolso: 5% (eventos subjetivos sem fontes claras)",
                1, 6, 11, 0.6,
                font_size=18, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = \
        "Sistema de 4 camadas: APIs fixas, IA com dupla verificação, retry automático, reembolso. " \
        "Taxa: 87% resolução automática, 5% reembolso (eventos subjetivos). " \
        "APIs: Banco Central, Yahoo Finance, CoinGecko, TSE. " \
        "IA: 2 checks independentes, confiança ≥85%."
    
    # ========== SLIDE 9: ESTADO ATUAL ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Estado Atual (Maio/2026)", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    # Esquerda
    add_textbox(slide, "O que somos hoje", 1, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_hoje = [
        "• Lançamos utilizando Zafes como moeda virtual",
        "• Sistema de resolução em 4 camadas (87% sucesso)",
        "• Eventos reais rodando na Liga e Econômico"
    ]
    add_multiline_text(slide, items_hoje, 1, 2.6, 5.5, 1.5,
                       font_size=16, color=WHITE)
    
    add_textbox(slide, "Objetivo Imediato: Alcançar 1.000 usuários ativos na fase de testes.",
                1, 4.2, 5.5, 0.6,
                font_size=16, bold=True, color=YELLOW)
    
    # Direita
    add_textbox(slide, "O que muda com 1000 usuários", 7, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_muda = [
        "• Aumentar volume de palpites na Liga",
        "• Validar modelo de competição social",
        "• Migrar de beta para lançamento oficial",
        "• Primeiro concurso com prêmio real via PIX",
        "• Atrair patrocinadores corporativos"
    ]
    add_multiline_text(slide, items_muda, 7, 2.6, 5.5, 2,
                       font_size=16, color=WHITE)
    
    add_textbox(slide, "Zafes: A moeda virtual que gera prêmios reais através de comprovação de habilidade analítica.",
                1, 6.2, 11, 0.6,
                font_size=18, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = \
        "Estado atual: maio/2026, usando Zafes como moeda virtual. " \
        "Sistema 4 camadas: 87% sucesso na resolução automática. " \
        "Objetivo: 1.000 usuários ativos na fase beta. " \
        "Com 1000 usuários: migramos para lançamento oficial."
    
    # ========== SLIDE 10: CURTO PRAZO ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Curto Prazo (1-3 meses)", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    # Esquerda
    add_textbox(slide, "Estruturação Jurídica e Digital", 1, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_curto_1 = [
        "• Criar CNPJ: Constituir empresa formalmente",
        "• Comprar Domínio: zafe.com.br (migração de zafe-rho.vercel.app)",
        "• Email Corporativo: @zafe.com.br",
        "• Conta Bancária PJ: Receber assinaturas e pagar prêmios"
    ]
    add_multiline_text(slide, items_curto_1, 1, 2.6, 5.5, 2,
                       font_size=16, color=WHITE)
    
    # Direita
    add_textbox(slide, "Marketing e Lançamento", 7, 2, 5, 0.5,
                font_size=24, bold=True, color=GREEN)
    
    items_curto_2 = [
        "• Criar Redes Sociais: Instagram, Twitter, TikTok, LinkedIn",
        "• Campanha de Advertising: Ads 18-35 anos",
        "• Lançar Oficialmente: Rebrand Zafes para Zafe",
        "• Alcançar 1.000 usuários: Tração orgânica + paga",
        "• App Mobile PWA: Notificações nativas"
    ]
    add_multiline_text(slide, items_curto_2, 7, 2.6, 5.5, 2.5,
                       font_size=16, color=WHITE)
    
    slide.notes_slide.notes_text_frame.text = \
        "Curto prazo (1-3 meses): Criar CNPJ, comprar domínio zafe.com.br. " \
        "Email corporativo @zafe.com.br, conta bancária PJ. " \
        "Marketing: redes sociais, ads 18-35 anos. " \
        "Rebrand para Zafe, meta 1.000 usuários, app PWA."
    
    # ========== SLIDE 11: MEDIO LONGO PRAZO ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "Roadmap: Médio e Longo Prazo", 1, 0.5, 11, 1,
                font_size=44, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    # Médio prazo
    add_textbox(slide, "Médio Prazo (3-6 meses)", 1, 2, 3.5, 0.5,
                font_size=20, bold=True, color=GREEN)
    
    items_medio = [
        "• Zafe Premium (R$ 19,90/mês)",
        "• Parcerias de patrocínio corporativo",
        "• Concurso Trimestral R$ 25k",
        "• Gamificação (badges, níveis, streaks)",
        "• Recursos sociais (seguir previsores)"
    ]
    add_multiline_text(slide, items_medio, 1, 2.6, 3.8, 2,
                       font_size=14, color=WHITE)
    
    # Longo prazo
    add_textbox(slide, "Longo Prazo (6-12 meses)", 5, 2, 3.5, 0.5,
                font_size=20, bold=True, color=GREEN)
    
    items_longo = [
        "• Troca de Zafes: 10.000 Z$ = 10% desconto",
        "• Zafes viram moeda de fidelidade",
        "• Converter Z$ em descontos reais (parcerias)",
        "• Venda de dados B2B (sabedoria das multidões)",
        "• Premium escalado nacionalmente"
    ]
    add_multiline_text(slide, items_longo, 5, 2.6, 3.8, 2,
                       font_size=14, color=WHITE)
    
    # Visão expansão
    add_textbox(slide, "Visão de Expansão", 9, 2, 3.5, 0.5,
                font_size=20, bold=True, color=GREEN)
    
    items_exp = [
        "• Consolidação da marca nacionalmente",
        "• Liderança em liga de previsões no Brasil",
        "• Parcerias com varejistas",
        "• Expansão para 100k+ usuários"
    ]
    add_multiline_text(slide, items_exp, 9, 2.6, 3.8, 1.6,
                       font_size=14, color=WHITE)
    
    slide.notes_slide.notes_text_frame.text = \
        "Médio prazo: Zafe Premium R$ 19,90/mês, concurso trimestral R$ 25k. " \
        "Longo prazo: Zafes viram moeda de fidelidade, conversão em descontos. " \
        "Visão: liderança nacional, 100k+ usuários. " \
        "Venda de dados B2B (sabedoria das multidões)."
    
    # ========== SLIDE 12: FECHAMENTO ==========
    slide = add_slide(prs)
    set_slide_background(slide, BLACK)
    
    add_textbox(slide, "ZAFE", 1, 1, 11, 1.5,
                font_size=80, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "A Liga de Previsões do Brasil", 1, 2.5, 11, 0.8,
                font_size=32, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    
    items_fech = [
        "• 100% Legal (Lei 5.768/71 + CMN 5.298)",
        "• Sem apostas, sem depósitos de dinheiro real",
        "• Prêmios reais no concurso via PIX",
        "• Tecnologia de ponta (4 camadas de resolução)",
        "• Zafes: a moeda virtual que vale descontos reais"
    ]
    add_multiline_text(slide, items_fech, 2, 3.5, 9, 2,
                       font_size=18, color=WHITE,
                       alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "Zafe: A liga onde você compete prevendo o que vai acontecer.",
                1, 5.8, 11, 0.6,
                font_size=20, bold=True, color=GREEN,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "zafe-rho.vercel.app → em breve zafe.com.br",
                1, 6.5, 11, 0.4,
                font_size=16, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, "FEITO POR LUC SAPOZNIK E LORENZO FRAGALI",
                1, 7, 11, 0.5,
                font_size=18, bold=True, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    
    slide.notes_slide.notes_text_frame.text = \
        "Fechamento: Zafe é a liga de previsões do Brasil (100% legal). " \
        "Lei 5.768/71 + CMN 5.298: sem apostas, sem depósitos reais. " \
        "Prêmios reais via PIX no concurso, 4 camadas de resolução. " \
        "Obrigado pela atenção. zafe-rho.vercel.app → zafe.com.br."
    
    # Salvar apresentação
    output_path = "/Users/mac/Downloads/zafe/apresentacao_zafe.pptx"
    prs.save(output_path)
    print(f"Apresentação salva em: {output_path}")
    return output_path

if __name__ == "__main__":
    main()
